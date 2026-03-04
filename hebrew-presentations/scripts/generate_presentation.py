#!/usr/bin/env python3
"""Generate Hebrew RTL PowerPoint presentations.

Creates properly formatted Hebrew presentations with RTL support,
Israeli business conventions, and common slide templates.

Usage:
    python generate_presentation.py --template pitch --output pitch.pptx
    python generate_presentation.py --template report --output report.pptx
    python generate_presentation.py --help

Requirements:
    pip install python-pptx
"""

import argparse
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


def set_rtl(paragraph):
    """Set paragraph direction to RTL."""
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    pPr.set("algn", "r")
    pPr.set("rtl", "1")


def add_hebrew_text(text_frame, text, font_size=18, bold=False, font_name="Heebo"):
    """Add RTL Hebrew text to a text frame."""
    paragraph = text_frame.paragraphs[0]
    set_rtl(paragraph)
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn("a:cs"), font_name)
    return paragraph


def create_pitch_deck(output_path):
    """Create an Israeli startup pitch deck template."""
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    slides_content = [
        ("שם הסטארטאפ", "משפט תיאור קצר של המוצר"),
        ("הבעיה", "מהי הבעיה שאנחנו פותרים?"),
        ("הפתרון", "כיצד המוצר שלנו פותר את הבעיה"),
        ("גודל השוק", "TAM / SAM / SOM"),
        ("המודל העסקי", "כיצד אנחנו מייצרים הכנסות"),
        ("Traction", "מספרים ואבני דרך"),
        ("הצוות", "מי אנחנו"),
        ("התחרות", "כיצד אנחנו שונים"),
        ("פיננסים", "תחזית הכנסות"),
        ("הבקשה", "כמה אנחנו מגייסים ולמה"),
    ]

    for title, subtitle in slides_content:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        add_hebrew_text(title_shape.text_frame, title, font_size=36, bold=True)

        body_shape = slide.placeholders[1]
        add_hebrew_text(body_shape.text_frame, subtitle, font_size=24)

    prs.save(output_path)
    print(f"Created pitch deck: {output_path}")
    print(f"  {len(slides_content)} slides")


def create_report(output_path):
    """Create a Hebrew business report template."""
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    slides_content = [
        ("דוח [נושא]", "[תאריך]"),
        ("תקציר מנהלים", "נקודות עיקריות"),
        ("רקע", "הקשר ומטרות"),
        ("ממצאים", "תוצאות עיקריות"),
        ("ניתוח", "משמעות הנתונים"),
        ("המלצות", "צעדים מוצעים"),
        ("סיכום", "צעדים הבאים ולוח זמנים"),
    ]

    for title, subtitle in slides_content:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        add_hebrew_text(title_shape.text_frame, title, font_size=36, bold=True)

        body_shape = slide.placeholders[1]
        add_hebrew_text(body_shape.text_frame, subtitle, font_size=24)

    prs.save(output_path)
    print(f"Created report template: {output_path}")
    print(f"  {len(slides_content)} slides")


def main():
    parser = argparse.ArgumentParser(
        description="Generate Hebrew RTL PowerPoint presentations"
    )
    parser.add_argument(
        "--template", choices=["pitch", "report"],
        default="pitch",
        help="Presentation template (default: pitch)"
    )
    parser.add_argument(
        "--output", default="presentation.pptx",
        help="Output file path (default: presentation.pptx)"
    )
    args = parser.parse_args()

    if args.template == "pitch":
        create_pitch_deck(args.output)
    elif args.template == "report":
        create_report(args.output)


if __name__ == "__main__":
    main()
